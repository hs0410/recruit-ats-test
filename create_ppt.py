#!/usr/bin/env python3
"""채용 ATS 서비스 설계서 PPT 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 색상 팔레트
PRIMARY = RGBColor(0x31, 0x82, 0xF6)
PRIMARY_DARK = RGBColor(0x1A, 0x56, 0xDB)
INDIGO = RGBColor(0x63, 0x66, 0xF1)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
GREEN = RGBColor(0x00, 0xC4, 0x71)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xF0, 0x44, 0x52)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x19, 0x1F, 0x28)
GRAY = RGBColor(0x8B, 0x95, 0xA1)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_gradient_bg(slide, color1=PRIMARY, color2=INDIGO):
    """그라데이션 배경 (단색 대체)"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color1


def add_light_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG


def add_white_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_text(slide, left, top, width, height, text, font_size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# ============================================================
# SLIDE 1: 표지
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_gradient_bg(slide, PRIMARY)

# 장식 원
add_circle(slide, Inches(-1), Inches(-1), Inches(4), RGBColor(0x3D, 0x8E, 0xF7))
add_circle(slide, Inches(10), Inches(4), Inches(5), RGBColor(0x4A, 0x75, 0xEB))

add_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
         '채용 ATS', 52, True, WHITE, PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(2.9), Inches(10), Inches(0.8),
         'Applicant Tracking System', 28, False, RGBColor(0xBB, 0xD5, 0xFC), PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.6),
         '채용의 모든 과정을 한 곳에서 관리하는 올인원 채용 플랫폼', 20, False, WHITE, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.5),
         '서비스 설계서  |  2026.04', 16, False, RGBColor(0x99, 0xBE, 0xF5), PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: 목차
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.8),
         'Contents', 36, True, PRIMARY)
add_text(slide, Inches(1), Inches(1.3), Inches(11), Inches(0.5),
         '목차', 16, False, GRAY)

toc_items = [
    ('01', '문제 정의', '채용 담당자가 겪는 핵심 문제'),
    ('02', '서비스 개요', '솔루션과 핵심 가치'),
    ('03', '핵심 기능', '10가지 주요 기능'),
    ('04', '페이지 구조', '전체 IA 및 네비게이션'),
    ('05', '기술 스택', '프레임워크, DB, 배포'),
    ('06', 'DB 스키마', '데이터 모델 설계'),
    ('07', 'HiringKit 통합', '기존 앱 기능 융합'),
    ('08', '향후 로드맵', '다음 단계 계획'),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(2.2) + Inches(i * 0.6)
    add_rounded_rect(slide, Inches(1), y, Inches(11), Inches(0.5), LIGHT_BG)
    add_text(slide, Inches(1.3), y + Pt(4), Inches(0.8), Inches(0.4),
             num, 16, True, PRIMARY)
    add_text(slide, Inches(2.2), y + Pt(4), Inches(3), Inches(0.4),
             title, 16, True, BLACK)
    add_text(slide, Inches(5.5), y + Pt(4), Inches(6), Inches(0.4),
             desc, 14, False, GRAY)


# ============================================================
# SLIDE 3: 문제 정의
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(4), Inches(0.7),
         '01  문제 정의', 32, True, PRIMARY)
add_text(slide, Inches(1), Inches(1.3), Inches(11), Inches(0.5),
         '채용 담당자가 매일 겪는 핵심 문제 3가지', 16, False, GRAY)

problems = [
    ('채용 브랜딩의 부재', '획일적인 JD로 우수 인재 확보 실패\n급여, 성장기회, 문화가 제대로 전달 안됨\n경쟁사 대비 매력도가 떨어지는 채용 공고', RED),
    ('리드타임 관리 불가', '후보자가 어느 단계에 얼마나 머물렀는지 파악 불가\n면접관 일정 조율에 과도한 시간 소모\n채용 병목 구간을 찾을 수 없음', ORANGE),
    ('휴먼에러 & 법적 리스크', '금지 질문(출신지역, 혼인여부) 실수 발생\n불합격 통보 누락, 서류 반환 미이행\n면접 평가 기록이 남지 않아 분쟁 위험', PURPLE),
]

for i, (title, desc, color) in enumerate(problems):
    x = Inches(1) + Inches(i * 3.8)
    card = add_rounded_rect(slide, x, Inches(2.3), Inches(3.5), Inches(4.2), CARD_BG, RGBColor(0xE5, 0xE8, 0xEB))
    add_circle(slide, x + Inches(1.3), Inches(2.6), Inches(0.8), color)
    add_text(slide, x + Inches(1.45), Inches(2.7), Inches(0.8), Inches(0.7),
             str(i + 1), 28, True, WHITE, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.3), Inches(3.6), Inches(2.9), Inches(0.5),
             title, 18, True, BLACK, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.3), Inches(4.2), Inches(2.9), Inches(2),
             desc, 12, False, GRAY, PP_ALIGN.LEFT)


# ============================================================
# SLIDE 4: 서비스 개요
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(4), Inches(0.7),
         '02  서비스 개요', 32, True, PRIMARY)

add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
         '채용 ATS는 채용 담당자와 면접관 모두를 위한\n올인원 채용 관리 플랫폼입니다.', 22, True, BLACK, PP_ALIGN.LEFT)

values = [
    ('후보자 파이프라인', '칸반보드로 전체 채용 현황을\n한눈에 파악', PRIMARY),
    ('AI 공고 작성', '5단계 위자드로\n전문적인 JD 자동 생성', INDIGO),
    ('면접관 가이드', '구조화 면접, 채용절차법\n교육까지 원스톱', PURPLE),
    ('데이터 기반', '리드타임, 병목 구간\n실시간 모니터링', GREEN),
]

for i, (title, desc, color) in enumerate(values):
    x = Inches(1) + Inches(i * 3)
    card = add_rounded_rect(slide, x, Inches(3.5), Inches(2.7), Inches(3), CARD_BG, RGBColor(0xE5, 0xE8, 0xEB))
    # 컬러 바
    bar = add_rounded_rect(slide, x, Inches(3.5), Inches(2.7), Inches(0.08), color)
    add_text(slide, x + Inches(0.3), Inches(3.9), Inches(2.2), Inches(0.5),
             title, 16, True, BLACK)
    add_text(slide, x + Inches(0.3), Inches(4.5), Inches(2.2), Inches(1.5),
             desc, 13, False, GRAY)


# ============================================================
# SLIDE 5: 대상 사용자
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.7),
         '02  대상 사용자', 32, True, PRIMARY)

# 채용 담당자
card1 = add_rounded_rect(slide, Inches(1), Inches(2), Inches(5.3), Inches(4.5), CARD_BG, RGBColor(0xE5, 0xE8, 0xEB))
add_rounded_rect(slide, Inches(1), Inches(2), Inches(5.3), Inches(0.08), PRIMARY)
add_text(slide, Inches(1.4), Inches(2.4), Inches(4.5), Inches(0.5),
         '채용 담당자 (HR)', 20, True, PRIMARY)
add_text(slide, Inches(1.4), Inches(3.1), Inches(4.5), Inches(3),
         '- 채용 공고 작성 및 관리\n- 후보자 파이프라인 운영\n- 면접 일정 조율\n- 리드타임 모니터링\n- 채용 데이터 분석\n- 채용절차법 준수 관리', 14, False, GRAY)

# 면접관
card2 = add_rounded_rect(slide, Inches(7), Inches(2), Inches(5.3), Inches(4.5), CARD_BG, RGBColor(0xE5, 0xE8, 0xEB))
add_rounded_rect(slide, Inches(7), Inches(2), Inches(5.3), Inches(0.08), PURPLE)
add_text(slide, Inches(7.4), Inches(2.4), Inches(4.5), Inches(0.5),
         '현업 면접관 (팀장)', 20, True, PURPLE)
add_text(slide, Inches(7.4), Inches(3.1), Inches(4.5), Inches(3),
         '- 면접 가이드 참고\n- STAR 기반 구조화 면접\n- 평가 기준표로 객관적 평가\n- 면접 전 체크리스트 확인\n- 무의식적 편향 교육\n- 면접 결과 기록', 14, False, GRAY)


# ============================================================
# SLIDE 6: 핵심 기능 (1/2)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.7),
         '03  핵심 기능 (1/2)', 32, True, PRIMARY)

features1 = [
    ('대시보드', '통계 카드, 파이프라인 현황,\n채용 팁 캐러셀, 빠른 이동', PRIMARY),
    ('칸반 파이프라인', '드래그&드롭 단계 이동,\n서류접수→1차→2차→처우→입사확정', INDIGO),
    ('후보자 관리', '목록/상세/등록, 검색 필터,\n타임라인, 메모, 리드타임 추적', PURPLE),
    ('채용 공고', '공고 목록/상세/편집,\n상태 관리 (진행중/마감)', GREEN),
    ('JD 위자드', '5단계 AI 공고 작성,\n11개 부서별 내부 템플릿', RGBColor(0xEC, 0x48, 0x99)),
]

for i, (title, desc, color) in enumerate(features1):
    col = i % 3
    row = i // 3
    x = Inches(1) + Inches(col * 3.8)
    y = Inches(1.8) + Inches(row * 2.8)
    card = add_rounded_rect(slide, x, y, Inches(3.5), Inches(2.4), CARD_BG, RGBColor(0xE5, 0xE8, 0xEB))
    add_rounded_rect(slide, x, y, Inches(0.08), Inches(2.4), color)
    add_text(slide, x + Inches(0.4), y + Inches(0.3), Inches(2.8), Inches(0.4),
             title, 16, True, BLACK)
    add_text(slide, x + Inches(0.4), y + Inches(0.9), Inches(2.8), Inches(1.2),
             desc, 12, False, GRAY)


# ============================================================
# SLIDE 7: 핵심 기능 (2/2)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.7),
         '03  핵심 기능 (2/2)', 32, True, PRIMARY)

features2 = [
    ('면접 일정', '주간 캘린더 뷰,\n일정 등록 모달, 면접관/장소 관리', PRIMARY),
    ('면접 가이드', '행동/기술/컬처핏 3탭,\n15개 질문 + STAR 프레임워크', INDIGO),
    ('면접관 교육', '구조화면접, 무의식편향,\n채용절차법, BEI 4개 모듈', PURPLE),
    ('평가 기준', '포지션별 역량 루브릭,\n1~5점 인터랙티브 평가표', GREEN),
    ('채용 인사이트', '트렌드/법률/AI 기사 8개,\n채용팁 16개, 2026 키워드', ORANGE),
]

for i, (title, desc, color) in enumerate(features2):
    col = i % 3
    row = i // 3
    x = Inches(1) + Inches(col * 3.8)
    y = Inches(1.8) + Inches(row * 2.8)
    card = add_rounded_rect(slide, x, y, Inches(3.5), Inches(2.4), CARD_BG, RGBColor(0xE5, 0xE8, 0xEB))
    add_rounded_rect(slide, x, y, Inches(0.08), Inches(2.4), color)
    add_text(slide, x + Inches(0.4), y + Inches(0.3), Inches(2.8), Inches(0.4),
             title, 16, True, BLACK)
    add_text(slide, x + Inches(0.4), y + Inches(0.9), Inches(2.8), Inches(1.2),
             desc, 12, False, GRAY)


# ============================================================
# SLIDE 8: 페이지 구조 (IA)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.7),
         '04  페이지 구조 (IA)', 32, True, PRIMARY)
add_text(slide, Inches(1), Inches(1.3), Inches(11), Inches(0.5),
         '총 15개 페이지 · 3개 섹션', 16, False, GRAY)

# 채용 관리 섹션
add_rounded_rect(slide, Inches(0.8), Inches(2), Inches(4), Inches(4.8), RGBColor(0xE8, 0xF3, 0xFF))
add_text(slide, Inches(1.1), Inches(2.2), Inches(3.5), Inches(0.4),
         '채용 관리', 16, True, PRIMARY)

mgmt_pages = ['/dashboard  대시보드', '/pipeline  파이프라인', '/candidates  후보자 목록',
              '/candidates/[id]  후보자 상세', '/candidates/new  후보자 등록',
              '/jobs  채용 공고', '/jobs/[id]  공고 상세', '/jobs/new  JD 위자드',
              '/schedule  면접 일정', '/settings  설정']
for i, page in enumerate(mgmt_pages):
    add_text(slide, Inches(1.1), Inches(2.8) + Inches(i * 0.37), Inches(3.5), Inches(0.35),
             page, 11, False, BLACK)

# 채용 도구 섹션
add_rounded_rect(slide, Inches(5.1), Inches(2), Inches(3.5), Inches(4.8), RGBColor(0xF3, 0xE8, 0xFF))
add_text(slide, Inches(5.4), Inches(2.2), Inches(3), Inches(0.4),
         '채용 도구 (HiringKit)', 16, True, PURPLE)

tool_pages = ['/guide  면접 가이드', '/guide/training  면접관 교육', '/guide/checklist  면접 전 체크리스트',
              '/evaluation  평가 기준 생성기', '/insights  채용 인사이트']
for i, page in enumerate(tool_pages):
    add_text(slide, Inches(5.4), Inches(2.8) + Inches(i * 0.37), Inches(3), Inches(0.35),
             page, 11, False, BLACK)

# 사이드바 구조
add_rounded_rect(slide, Inches(8.9), Inches(2), Inches(3.5), Inches(4.8), RGBColor(0xE8, 0xFA, 0xF0))
add_text(slide, Inches(9.2), Inches(2.2), Inches(3), Inches(0.4),
         '사이드바 네비게이션', 16, True, GREEN)

nav_items = ['대시보드', '파이프라인', '후보자', '채용 공고', '면접 일정',
             '── 채용 도구 ──', 'JD 작성', '면접 가이드', '평가 기준', '인사이트',
             '── ──', '설정']
for i, item in enumerate(nav_items):
    color = GRAY if '──' in item else BLACK
    add_text(slide, Inches(9.2), Inches(2.8) + Inches(i * 0.33), Inches(3), Inches(0.3),
             item, 11, False, color)


# ============================================================
# SLIDE 9: 기술 스택
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.7),
         '05  기술 스택', 32, True, PRIMARY)

stack = [
    ('프레임워크', 'Next.js 16 (App Router)', 'React 기반, 프론트+백엔드 통합', PRIMARY),
    ('언어', 'TypeScript', '타입 안정성, 개발 생산성', INDIGO),
    ('DB', 'Supabase (PostgreSQL)', '무료 플랜, 인증 내장, 실시간', GREEN),
    ('인증', 'Supabase Auth', '이메일/비밀번호 로그인', PURPLE),
    ('스타일링', 'Tailwind CSS', '유틸리티 기반, 빠른 개발', ORANGE),
    ('배포', 'Vercel', 'Next.js 최적, 무료 플랜', RED),
]

for i, (category, tech, reason, color) in enumerate(stack):
    col = i % 3
    row = i // 3
    x = Inches(1) + Inches(col * 3.8)
    y = Inches(1.8) + Inches(row * 2.6)
    card = add_rounded_rect(slide, x, y, Inches(3.5), Inches(2.2), CARD_BG, RGBColor(0xE5, 0xE8, 0xEB))
    add_text(slide, x + Inches(0.3), y + Inches(0.3), Inches(3), Inches(0.3),
             category, 12, True, color)
    add_text(slide, x + Inches(0.3), y + Inches(0.7), Inches(3), Inches(0.4),
             tech, 18, True, BLACK)
    add_text(slide, x + Inches(0.3), y + Inches(1.3), Inches(3), Inches(0.6),
             reason, 12, False, GRAY)

add_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         '총 비용: 0원 (Supabase 무료 + Vercel 무료)', 16, True, GREEN, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 10: DB 스키마
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.7),
         '06  DB 스키마', 32, True, PRIMARY)

tables = [
    ('profiles', 'id, email, name, role, created_at', '사용자 (관리자)'),
    ('jobs', 'id, title, department, description,\nstatus, created_at', '채용 공고'),
    ('candidates', 'id, name, phone, email, job_id,\nstage, applied_at, resume_url, memo', '후보자'),
    ('interviews', 'id, candidate_id, stage,\nscheduled_at, interviewer, location', '면접 일정'),
    ('stage_history', 'id, candidate_id, from_stage,\nto_stage, changed_at, changed_by', '단계 변경 이력'),
    ('evaluations', 'id, interview_id, candidate_id,\nrating, comment, created_by', '면접 평가'),
]

for i, (name, fields, desc) in enumerate(tables):
    col = i % 3
    row = i // 3
    x = Inches(1) + Inches(col * 3.8)
    y = Inches(1.8) + Inches(row * 2.7)
    card = add_rounded_rect(slide, x, y, Inches(3.5), Inches(2.3), CARD_BG, RGBColor(0xE5, 0xE8, 0xEB))
    add_rounded_rect(slide, x, y, Inches(3.5), Inches(0.06), PRIMARY)
    add_text(slide, x + Inches(0.3), y + Inches(0.2), Inches(2), Inches(0.4),
             name, 16, True, PRIMARY)
    add_text(slide, x + Inches(0.3), y + Inches(0.6), Inches(2.9), Inches(0.3),
             desc, 11, False, GRAY)
    add_text(slide, x + Inches(0.3), y + Inches(1), Inches(2.9), Inches(1),
             fields, 10, False, BLACK)


# ============================================================
# SLIDE 11: HiringKit 통합
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(8), Inches(0.7),
         '07  HiringKit 통합', 32, True, PRIMARY)
add_text(slide, Inches(1), Inches(1.3), Inches(11), Inches(0.5),
         '기존 1workplace(HiringKit) 앱의 5개 핵심 기능을 ATS에 통합', 16, False, GRAY)

integrations = [
    ('JD 위자드', '5단계 공고 작성 + AI 생성\n부서별 템플릿 11개', '/jobs/new', PRIMARY),
    ('면접 가이드', '행동/기술/컬처핏 질문\nSTAR 프레임워크', '/guide', INDIGO),
    ('면접관 교육', '구조화면접, 무의식편향\n채용절차법, BEI', '/guide/training', PURPLE),
    ('평가 기준', '포지션별 역량 루브릭\n인터랙티브 점수 입력', '/evaluation', GREEN),
    ('인사이트', '채용 트렌드/법률/AI 기사\n16개 채용 팁', '/insights', ORANGE),
]

for i, (title, desc, route, color) in enumerate(integrations):
    x = Inches(0.5) + Inches(i * 2.5)
    card = add_rounded_rect(slide, x, Inches(2.2), Inches(2.3), Inches(4), CARD_BG, RGBColor(0xE5, 0xE8, 0xEB))
    add_rounded_rect(slide, x, Inches(2.2), Inches(2.3), Inches(0.06), color)
    add_text(slide, x + Inches(0.2), Inches(2.6), Inches(1.9), Inches(0.4),
             title, 15, True, color)
    add_text(slide, x + Inches(0.2), Inches(3.2), Inches(1.9), Inches(1.2),
             desc, 11, False, GRAY)
    add_text(slide, x + Inches(0.2), Inches(5.5), Inches(1.9), Inches(0.3),
             route, 9, False, GRAY)

# 화살표 표시 대체
add_text(slide, Inches(1), Inches(6.6), Inches(11), Inches(0.5),
         'HiringKit의 오렌지(#F04E23) → ATS 블루(#3182F6) 계열로 디자인 통일', 13, False, GRAY, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 12: JD 위자드 상세
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(8), Inches(0.7),
         '03  JD 위자드 상세', 32, True, PRIMARY)
add_text(slide, Inches(1), Inches(1.3), Inches(11), Inches(0.5),
         '5단계로 완성하는 AI 채용 공고', 16, False, GRAY)

jd_steps = [
    ('Step 1\nWHY', '채용 목적\n(브랜딩 vs 효율)\n타겟 레벨 선택', PRIMARY),
    ('Step 2\nWHERE', '채널 선택\n(온라인/오프라인\n/하이브리드)', INDIGO),
    ('Step 3\nHOW', '상세 정보 입력\n제목, 부서, 업무\n자격, 연봉, 복리후생', PURPLE),
    ('Step 4\n미리보기', '부서별 내부 템플릿\n라이브러리 참고\n(11개 샘플)', GREEN),
    ('Step 5\n완성', 'AI 자동 생성\n복사 / 재생성\n공고 등록', RGBColor(0xEC, 0x48, 0x99)),
]

for i, (step, desc, color) in enumerate(jd_steps):
    x = Inches(0.5) + Inches(i * 2.5)
    # 원형 스텝
    circle = add_circle(slide, x + Inches(0.65), Inches(2.2), Inches(1), color)
    add_text(slide, x + Inches(0.65), Inches(2.35), Inches(1), Inches(0.8),
             step, 12, True, WHITE, PP_ALIGN.CENTER)
    # 설명
    add_text(slide, x + Inches(0.15), Inches(3.5), Inches(2), Inches(1.5),
             desc, 11, False, GRAY, PP_ALIGN.CENTER)
    # 화살표
    if i < 4:
        add_text(slide, x + Inches(1.8), Inches(2.5), Inches(0.5), Inches(0.5),
                 '→', 20, True, RGBColor(0xD0, 0xD5, 0xDD), PP_ALIGN.CENTER)

# 템플릿 부서
add_rounded_rect(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.5), LIGHT_BG)
add_text(slide, Inches(1.2), Inches(5.5), Inches(10), Inches(0.4),
         '내부 템플릿 라이브러리 (11개)', 14, True, BLACK)
depts = ['개발 (3개)', '영업 (2개)', 'HR (1개)', '품질 (2개)', '마케팅 (3개)']
for i, dept in enumerate(depts):
    add_rounded_rect(slide, Inches(1.2) + Inches(i * 2.2), Inches(6), Inches(2), Inches(0.5), CARD_BG, RGBColor(0xE5, 0xE8, 0xEB))
    add_text(slide, Inches(1.2) + Inches(i * 2.2), Inches(6.05), Inches(2), Inches(0.5),
             dept, 12, True, PRIMARY, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 13: 디자인 시스템
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.7),
         '디자인 시스템', 32, True, PRIMARY)
add_text(slide, Inches(1), Inches(1.3), Inches(11), Inches(0.5),
         'Toss / Notion 스타일 — 모던 & 클린', 16, False, GRAY)

# 컬러 팔레트
colors_display = [
    ('Primary', '#3182F6', PRIMARY),
    ('Indigo', '#6366F1', INDIGO),
    ('Purple', '#8B5CF6', PURPLE),
    ('Success', '#00C471', GREEN),
    ('Warning', '#F59E0B', ORANGE),
    ('Danger', '#F04452', RED),
]

add_text(slide, Inches(1), Inches(2), Inches(3), Inches(0.4),
         '컬러 팔레트', 16, True, BLACK)

for i, (name, hex_code, color) in enumerate(colors_display):
    x = Inches(1) + Inches(i * 1.8)
    add_rounded_rect(slide, x, Inches(2.6), Inches(1.5), Inches(0.9), color)
    add_text(slide, x, Inches(3.6), Inches(1.5), Inches(0.3),
             name, 11, True, BLACK, PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.9), Inches(1.5), Inches(0.3),
             hex_code, 10, False, GRAY, PP_ALIGN.CENTER)

# 디자인 원칙
add_text(slide, Inches(1), Inches(4.6), Inches(3), Inches(0.4),
         '디자인 원칙', 16, True, BLACK)

principles = [
    'Glassmorphism (반투명 글래스 효과)',
    '그라데이션 아이콘 & 배경',
    'Rounded 2XL (16px 라운드)',
    'Soft Shadow (은은한 그림자)',
    'Hover Lift (호버 시 살짝 떠오르는 효과)',
    'lucide-react 아이콘 통일',
]

for i, principle in enumerate(principles):
    add_text(slide, Inches(1.3), Inches(5.1) + Inches(i * 0.35), Inches(10), Inches(0.35),
             f'•  {principle}', 13, False, GRAY)


# ============================================================
# SLIDE 14: 향후 로드맵
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_bg(slide)

add_text(slide, Inches(1), Inches(0.6), Inches(5), Inches(0.7),
         '08  향후 로드맵', 32, True, PRIMARY)

phases = [
    ('현재', 'UI 프로토타입 완성\nMock 데이터 기반\n15개 페이지 구현', PRIMARY, True),
    ('Phase 1', 'Supabase 연동\n실제 데이터 CRUD\n인증 시스템 구현', INDIGO, False),
    ('Phase 2', '공개 채용 페이지\n지원 폼 → DB 자동 등록\n면접 리마인드 알림', PURPLE, False),
    ('Phase 3', 'Vercel 배포\n커스텀 도메인\n모바일 반응형', GREEN, False),
    ('Phase 4', 'AI 고도화\n실제 LLM API 연동\n이력서 파싱 자동화', ORANGE, False),
]

for i, (phase, desc, color, is_current) in enumerate(phases):
    x = Inches(0.5) + Inches(i * 2.5)
    y = Inches(2)
    card = add_rounded_rect(slide, x, y, Inches(2.3), Inches(4.5), CARD_BG, color if is_current else RGBColor(0xE5, 0xE8, 0xEB))
    if is_current:
        add_rounded_rect(slide, x + Inches(0.1), y + Inches(0.1), Inches(2.1), Inches(4.3), color)
        text_color = WHITE
    else:
        text_color = BLACK

    add_text(slide, x + Inches(0.3), y + Inches(0.4), Inches(1.8), Inches(0.4),
             phase, 18, True, text_color if is_current else color)
    add_text(slide, x + Inches(0.3), y + Inches(1.2), Inches(1.8), Inches(2.5),
             desc, 12, False, WHITE if is_current else GRAY)


# ============================================================
# SLIDE 15: 마무리
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, PRIMARY)

add_circle(slide, Inches(9), Inches(-1), Inches(5), RGBColor(0x4A, 0x75, 0xEB))
add_circle(slide, Inches(-2), Inches(4), Inches(4), RGBColor(0x3D, 0x8E, 0xF7))

add_text(slide, Inches(1.5), Inches(2), Inches(10), Inches(1),
         'Thank You', 48, True, WHITE, PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(1),
         '채용의 모든 과정을 한 곳에서.\n더 빠르고, 더 공정하고, 더 스마트한 채용.', 20, False, RGBColor(0xBB, 0xD5, 0xFC), PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
         'Powered by Claude Code  |  Next.js + Supabase + Vercel', 14, False, RGBColor(0x99, 0xBE, 0xF5), PP_ALIGN.CENTER)


# ============================================================
# 저장
# ============================================================
output_path = '/Users/hyunsoo/Documents/Claude/Scheduled/hr-english-buddy-daily/disney agent/ats-web/채용ATS_서비스설계서.pptx'
prs.save(output_path)
print(f'PPT 생성 완료: {output_path}')
print(f'총 {len(prs.slides)}장 슬라이드')
